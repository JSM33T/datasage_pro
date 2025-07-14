import os
import subprocess
import pickle
import faiss
import numpy as np
from pathlib import Path
from datetime import datetime
from fastapi import APIRouter, Body
from fastapi.responses import JSONResponse
from dotenv import load_dotenv
from pymongo import MongoClient
import openai
import fitz
import docx2txt
from pptx import Presentation
from nltk import sent_tokenize
import nltk
import tiktoken
encoding = tiktoken.get_encoding("cl100k_base")  # for text-embedding-ada-002

nltk.data.path.append("/home/azureadmin/nltk_data")

# More robust NLTK data handling
def ensure_nltk_data():
    """Ensure NLTK punkt tokenizer is available with fallback options."""
    try:
        # Try to find punkt_tab first (newer version)
        nltk.data.find('tokenizers/punkt_tab')
        return True
    except LookupError:
        try:
            # Try to download punkt_tab
            nltk.download('punkt_tab', quiet=True)
            return True
        except Exception:
            try:
                # Fallback to older punkt version
                nltk.data.find('tokenizers/punkt')
                return True
            except LookupError:
                try:
                    nltk.download('punkt', quiet=True)
                    return True
                except Exception:
                    return False

# Initialize NLTK data
nltk_available = ensure_nltk_data()

load_dotenv()

router = APIRouter()

# ENV + DB
RESOURCE_DIR = Path("./resources")
client = MongoClient(os.getenv("MONGO_URI"))
db = client[os.getenv("MONGO_DB")]  # type: ignore
docs = db["documents"]

openai_client = openai.OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

def get_embedding(text: str):
    response = openai_client.embeddings.create(
        model="text-embedding-ada-002",
        input=text
    )
    return response.data[0].embedding

def convert_ppt_to_pptx(file_path: Path) -> Path:
    """
    Converts a .ppt file to .pptx using LibreOffice CLI.
    Returns the new .pptx Path if successful, else raises Exception.
    """
    output_dir = file_path.parent
    cmd = [
        "libreoffice",
        "--headless",
        "--convert-to", "pptx",
        str(file_path),
        "--outdir", str(output_dir)
    ]
    result = subprocess.run(cmd, capture_output=True)
    if result.returncode != 0:
        raise Exception(f"LibreOffice conversion failed: {result.stderr.decode()}")
    new_file = output_dir / (file_path.stem + ".pptx")
    if not new_file.exists():
        raise Exception("Converted .pptx file not found after conversion.")
    return new_file

@router.post("/index")
async def index_document(data: dict = Body(...)):
    ids = data.get("doc_ids")
    if not ids or not isinstance(ids, list):
        return JSONResponse(status_code=400, content={"error": "doc_ids must be a list"})

    results = []
    for doc_id in ids:
        try:
            record = docs.find_one({"_id": doc_id})
            if not record:
                results.append({"id": doc_id, "status": "not_found"})
                continue

            doc_dir = RESOURCE_DIR / doc_id
            file_path = doc_dir / record["filename"]
            if not file_path.exists():
                results.append({"id": doc_id, "status": "file_missing"})
                continue

            # Extract text robustly
            ext = file_path.suffix.lower()
            text = ""
            extraction_error = None
            try:
                if ext == ".pdf":
                    try:
                        doc = fitz.open(str(file_path))
                        text = "\n".join(page.get_text() for page in doc)
                        if not text.strip():
                            extraction_error = "PDF text extraction empty (possibly scanned or encrypted)"
                    except Exception as e:
                        extraction_error = f"PDF extraction error: {str(e)}"

                elif ext in [".docx", ".doc"]:
                    try:
                        text = docx2txt.process(str(file_path))
                        if not text.strip():
                            extraction_error = "DOCX text extraction empty"
                    except Exception as e:
                        extraction_error = f"DOCX extraction error: {str(e)}"

                elif ext in [".ppt", ".pptx"]:
                    try:
                        # If .ppt, convert to .pptx first
                        if ext == ".ppt":
                            file_path = convert_ppt_to_pptx(file_path)

                        prs = Presentation(str(file_path))
                        slides_text = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text") and shape.text:
                                    slides_text.append(shape.text)
                        text = "\n".join(slides_text)
                        if not text.strip():
                            extraction_error = "PPTX text extraction empty"
                    except Exception as e:
                        extraction_error = f"PPTX extraction error: {str(e)}"

                else:
                    results.append({"id": doc_id, "status": f"unsupported_file: {ext}"})
                    continue
            except Exception as e:
                extraction_error = f"General extraction error: {str(e)}"

            if extraction_error:
                results.append({"id": doc_id, "status": "extract_error", "detail": extraction_error})
                continue

            # Chunking & Embedding
            chunks = chunk_text(text)
            # Remove empty/very short chunks
            chunks = [c for c in chunks if c and len(c.strip()) > 10]
            if not chunks:
                results.append({"id": doc_id, "status": "no_valid_chunks"})
                continue

            embeddings = []
            embedding_dim = None
            for chunk in chunks:
                try:
                    emb = get_embedding(chunk)
                    if embedding_dim is None:
                        embedding_dim = len(emb)
                    elif len(emb) != embedding_dim:
                        results.append({"id": doc_id, "status": "embedding_dim_mismatch"})
                        break
                    embeddings.append(emb)
                except Exception as e:
                    results.append({"id": doc_id, "status": "embedding_error", "detail": str(e)})
                    break
            else:
                if not embeddings:
                    results.append({"id": doc_id, "status": "no_embeddings"})
                    continue

                embedding = np.array(embeddings, dtype="float32")
                # Create FAISS index
                try:
                    dim = embedding.shape[1]
                    index = faiss.IndexFlatL2(dim)
                    index.add(embedding)
                except Exception as e:
                    results.append({"id": doc_id, "status": "faiss_error", "detail": str(e)})
                    continue

                # Atomic write: write to temp then move
                try:
                    faiss_tmp = doc_dir / f"{doc_id}.faiss.tmp"
                    faiss_final = doc_dir / f"{doc_id}.faiss"
                    faiss.write_index(index, str(faiss_tmp))
                    os.replace(faiss_tmp, faiss_final)
                    pkl_tmp = doc_dir / f"{doc_id}.pkl.tmp"
                    pkl_final = doc_dir / f"{doc_id}.pkl"
                    with open(pkl_tmp, "wb") as f:
                        pickle.dump({"text": chunks, "filename": record["filename"]}, f)
                    os.replace(pkl_tmp, pkl_final)
                except Exception as e:
                    # Clean up temp files
                    for f in [faiss_tmp, pkl_tmp]:
                        try:
                            if f.exists():
                                f.unlink()
                        except Exception:
                            pass
                    results.append({"id": doc_id, "status": "atomic_write_error", "detail": str(e)})
                    continue

                # Update DB only if all succeeded
                docs.update_one({"_id": doc_id}, {"$set": {"isIndexed": True}})
                results.append({"id": doc_id, "status": "indexed", "chunks": len(chunks)})

        except Exception as e:
            results.append({"id": doc_id, "status": f"error", "detail": str(e)})

    return {"results": results}

def chunk_text(text, max_tokens=200, overlap=50):
    sentences = sent_tokenize(text)
    chunks = []
    current_chunk = []
    current_tokens = 0

    for sentence in sentences:
        tokens = encoding.encode(sentence)
        token_count = len(tokens)

        if current_tokens + token_count > max_tokens:
            chunks.append(" ".join(current_chunk))
            # add overlap tokens from previous chunk
            overlap_tokens = encoding.encode(" ".join(current_chunk))[-overlap:]
            current_chunk = encoding.decode(overlap_tokens).split() + sentence.split()
            current_tokens = len(encoding.encode(" ".join(current_chunk)))
        else:
            current_chunk.append(sentence)
            current_tokens += token_count

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    return chunks


def chunk_text_old(text, max_tokens=200, overlap=50):
    """
    Chunk text into smaller segments with overlap.
    Falls back to simple splitting if NLTK is not available.
    """
    global nltk_available
    
    if nltk_available:
        try:
            sentences = sent_tokenize(text)
        except Exception:
            # Fallback to simple sentence splitting
            sentences = text.replace('!', '.').replace('?', '.').split('.')
            sentences = [s.strip() for s in sentences if s.strip()]
    else:
        # Simple sentence splitting fallback
        sentences = text.replace('!', '.').replace('?', '.').split('.')
        sentences = [s.strip() for s in sentences if s.strip()]
    
    chunks = []
    current_chunk = []
    current_len = 0

    for sentence in sentences:
        tokens = sentence.split()
        if current_len + len(tokens) > max_tokens:
            if current_chunk:  # Only append if current_chunk is not empty
                chunks.append(" ".join(current_chunk))
            # add overlap
            current_chunk = current_chunk[-overlap:] + tokens
            current_len = len(current_chunk)
        else:
            current_chunk += tokens
            current_len += len(tokens)

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    return chunks if chunks else [text]  # Return original text if no chunks created